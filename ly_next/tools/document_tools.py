from __future__ import annotations

import asyncio
from typing import Any

from ly_next.tools.base import ToolResult, tool
from ly_next.tools.export_paths import export_download_url, exports_dir, safe_export_filename


def _parse_sections(sections: Any) -> list[dict[str, Any]]:
    if not isinstance(sections, list):
        return []
    out: list[dict[str, Any]] = []
    for item in sections:
        if isinstance(item, str) and item.strip():
            out.append({"heading": "", "paragraphs": [item.strip()]})
            continue
        if not isinstance(item, dict):
            continue
        heading = str(item.get("heading") or item.get("title") or "").strip()
        raw = item.get("paragraphs") or item.get("content") or item.get("body")
        paras: list[str] = []
        if isinstance(raw, str) and raw.strip():
            paras = [raw.strip()]
        elif isinstance(raw, list):
            paras = [str(p).strip() for p in raw if str(p).strip()]
        if heading or paras:
            out.append({"heading": heading, "paragraphs": paras})
    return out


def _parse_slides(slides: Any) -> list[dict[str, Any]]:
    if not isinstance(slides, list):
        return []
    out: list[dict[str, Any]] = []
    for item in slides:
        if not isinstance(item, dict):
            continue
        title = str(item.get("title") or item.get("heading") or "").strip()
        raw = item.get("bullets") or item.get("points") or item.get("content")
        bullets: list[str] = []
        if isinstance(raw, str) and raw.strip():
            bullets = [raw.strip()]
        elif isinstance(raw, list):
            bullets = [str(b).strip() for b in raw if str(b).strip()]
        if title or bullets:
            out.append({"title": title, "bullets": bullets})
    return out


def _write_docx(path: str, *, title: str, sections: list[dict[str, Any]]) -> None:
    from docx import Document

    doc = Document()
    if title.strip():
        doc.add_heading(title.strip(), level=0)
    for sec in sections:
        if sec.get("heading"):
            doc.add_heading(str(sec["heading"]), level=1)
        for para in sec.get("paragraphs") or []:
            if str(para).strip():
                doc.add_paragraph(str(para).strip())
    doc.save(path)


def _write_xlsx(path: str, *, sheet_name: str, headers: list[str], rows: list[list[Any]]) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = (sheet_name or "Sheet1")[:31] or "Sheet1"
    if headers:
        ws.append([str(h) for h in headers])
    for row in rows:
        if isinstance(row, list):
            ws.append(row)
    wb.save(path)


def _write_pptx(path: str, *, title: str, slides: list[dict[str, Any]]) -> None:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    if title.strip():
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title.strip()
    for spec in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(spec.get("title") or "Slide")
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for i, bullet in enumerate(spec.get("bullets") or []):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = str(bullet)
            p.level = 0
            p.font.size = Pt(18)
    prs.save(path)


async def _export(*, ext: str, stem: str, writer) -> ToolResult:
    try:
        name = safe_export_filename(stem, ext)
        path = exports_dir() / name
        await asyncio.to_thread(writer, str(path))
        return ToolResult(
            success=True,
            result={
                "filename": name,
                "download_url": export_download_url(name),
                "format": ext,
            },
        )
    except ImportError as e:
        return ToolResult(success=False, error=str(e))
    except ValueError as e:
        return ToolResult(success=False, error=str(e))
    except Exception as e:
        return ToolResult(success=False, error=str(e))


async def generate_docx(
    title: str,
    sections: list[dict[str, Any]] | list[str] | None = None,
    filename: str = "",
) -> ToolResult:
    parsed = _parse_sections(sections)
    if not title.strip() and not parsed:
        return ToolResult(success=False, error="title or sections required")
    stem = filename or title or "document"

    def _run(path: str) -> None:
        try:
            _write_docx(path, title=title, sections=parsed)
        except ImportError as e:
            raise ImportError("install python-docx") from e

    return await _export(ext="docx", stem=stem, writer=_run)


async def generate_xlsx(
    headers: list[str],
    rows: list[list[Any]],
    sheet_name: str = "Sheet1",
    filename: str = "",
) -> ToolResult:
    if not headers and not rows:
        return ToolResult(success=False, error="headers or rows required")
    stem = filename or sheet_name or "spreadsheet"

    def _run(path: str) -> None:
        try:
            _write_xlsx(path, sheet_name=sheet_name, headers=headers, rows=rows)
        except ImportError as e:
            raise ImportError("install openpyxl") from e

    return await _export(ext="xlsx", stem=stem, writer=_run)


async def generate_pptx(
    slides: list[dict[str, Any]],
    title: str = "",
    filename: str = "",
) -> ToolResult:
    parsed = _parse_slides(slides)
    if not title.strip() and not parsed:
        return ToolResult(success=False, error="title or slides required")
    stem = filename or title or "presentation"

    def _run(path: str) -> None:
        try:
            _write_pptx(path, title=title, slides=parsed)
        except ImportError as e:
            raise ImportError("install python-pptx") from e

    return await _export(ext="pptx", stem=stem, writer=_run)


generate_docx_tool = tool(
    name="generate_docx",
    description=(
        "Create a Word .docx and return download_url. "
        "Pass title plus sections: [{heading, paragraphs: [text, ...]}]."
    ),
    category="general",
    parameters={
        "type": "object",
        "required": ["title"],
        "properties": {
            "title": {"type": "string", "description": "Document title."},
            "sections": {
                "type": "array",
                "description": "Sections with heading and paragraphs.",
                "items": {"type": "object"},
            },
            "filename": {"type": "string", "description": "Optional filename stem."},
        },
    },
)(generate_docx)

generate_xlsx_tool = tool(
    name="generate_xlsx",
    description=(
        "Create an Excel .xlsx and return download_url. "
        "Required: headers (string[]), rows (array of row arrays)."
    ),
    category="general",
    parameters={
        "type": "object",
        "required": ["headers", "rows"],
        "properties": {
            "headers": {"type": "array", "items": {"type": "string"}},
            "rows": {"type": "array", "items": {"type": "array"}},
            "sheet_name": {"type": "string", "default": "Sheet1"},
            "filename": {"type": "string"},
        },
    },
)(generate_xlsx)

generate_pptx_tool = tool(
    name="generate_pptx",
    description=(
        "Create a PowerPoint .pptx and return download_url. "
        "Pass slides: [{title, bullets: [text, ...]}]; optional cover title."
    ),
    category="general",
    parameters={
        "type": "object",
        "required": ["slides"],
        "properties": {
            "title": {"type": "string", "description": "Optional cover title."},
            "slides": {
                "type": "array",
                "description": "Slides with title and bullets.",
                "items": {"type": "object"},
            },
            "filename": {"type": "string"},
        },
    },
)(generate_pptx)
